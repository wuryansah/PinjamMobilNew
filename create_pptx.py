from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 212, 255)
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 215, 0)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(12)
    return slide

# SLIDE 1: Title
add_title_slide(prs, "PROSEDUR SISTEM", "PINJAM MOBIL")

# SLIDE 2: Gambaran Umum
add_content_slide(prs, "Gambaran Umum Sistem", [
    "Sistem Pinjam Mobil adalah aplikasi berbasis web untuk",
    "mengelola peminjaman kendaraan perusahaan.",
    "",
    "Fitur Utama:",
    "• Pengelolaan data kendaraan (armada)",
    "• Penerimaan dan persetujuan permintaan mobil",
    "• Penugasan driver dan kendaraan",
    "• Pelacakan penggunaan (start KM, end KM)",
    "• Monitor bahan bakar dan laporan"
])

# SLIDE 3: Jenis Pengguna
add_content_slide(prs, "Jenis Pengguna Sistem", [
    "• ADMIN - Kelola kendaraan, driver, pengguna, approve/request, generate laporan",
    "• MANAGER - Approve permintaan kendaraan dari bawahannya",
    "• DRIVER - Mengemudikan kendaraan, mulai/selesaikan perjalanan",
    "• KARYAWAN - Membuat permintaan peminjaman kendaraan"
])

# SLIDE 4: Login
add_content_slide(prs, "Langkah 1: Login Sistem", [
    "1. Buka halaman login sistem",
    "2. Masukkan email dan password",
    "3. Klik tombol 'Login'",
    "",
    "Catatan: Password default adalah 'password' untuk user baru.",
    "Admin harus mengganti setelah login pertama."
])

# SLIDE 5: Membuat Permintaan
add_content_slide(prs, "Langkah 2: Membuat Permintaan (Karyawan)", [
    "1. Klik menu 'Permintaan' atau 'Requests'",
    "2. Klik tombol '+ Buat Permintaan Baru'",
    "3. Isi formulir:",
    "   • Tujuan (destination)",
    "   • Keperluan (purpose)",
    "   • Tanggal/waktu mulai",
    "   • Tanggal/waktu selesai",
    "4. Klik 'Simpan' atau 'Submit'",
    "",
    "Status akan menjadi 'Pending' menunggu persetujuan manager."
])

# SLIDE 6: Persetujuan Manager
add_content_slide(prs, "Langkah 3: Persetujuan Manager", [
    "1. Manager menerima notifikasi permintaan baru",
    "2. Buka menu 'Permintaan'",
    "3. Lihat detail permintaan karyawan",
    "4. Pilih action:",
    "   • SETUJU - Kirim ke admin untuk penugasan",
    "   • TOLAK - Berikan alasan penolakan",
    "5. Klik tombol proses",
    "",
    "Jika user adalah Manager (self-request), status langsung",
    "menjadi 'Manager Approved'."
])

# SLIDE 7: Penugasaan Admin
add_content_slide(prs, "Langkah 4: Penugasaan Admin", [
    "1. Admin menerima notifikasi request yang di-approve manager",
    "2. Buka detail permintaan",
    "3. Pilih kendaraan yang tersedia",
    "4. Pilih driver yang akan bertugas",
    "5. Tambah catatan jika diperlukan",
    "6. Klik 'Setuju & Tugaskan'",
    "",
    "Alternatif: Admin bisa menolak dengan klik 'Tolak'."
])

# SLIDE 8: Memulai Perjalanan
add_content_slide(prs, "Langkah 5: Memulai Perjalanan (Driver)", [
    "1. Driver menerima notifikasi penugasan",
    "2. Buka menu 'Permintaan'",
    "3. Cari trip yang ditugaskan",
    "4. Klik 'Mulai Perjalanan'",
    "5. Input kilometer awal (start KM)",
    "6. Klik 'Simpan'",
    "",
    "Status berubah menjadi 'In Progress' (sedang berlangsung)."
])

# SLIDE 9: Menyelesaikan Perjalanan
add_content_slide(prs, "Langkah 6: Menyelesaikan Perjalanan (Driver)", [
    "1. Setelah sampai tujuan",
    "2. Klik 'Selesai Perjalanan'",
    "3. Input:",
    "   • Kilometer akhir (end KM)",
    "   • Catatan perjalanan",
    "4. Klik 'Simpan'",
    "5. Kendaraan kembali tersedia",
    "",
    "Status menjadi 'Completed' dan datausage seringkan untuk laporan."
])

# SLIDE 10: Pencatatan BBM
add_content_slide(prs, "Langkah 7: Pencatatan Bahan Bakar", [
    "1. Akses menu 'BBM' atau 'Fuel'",
    "2. Klik '+ Tambah Catatan BBM'",
    "3. Isi formulir:",
    "   • Pilih kendaraan",
    "   • Tanggal pengisisan",
    "   • Jumlah liter",
    "   • Total biaya",
    "   • Upload bukti nota (opsional)",
    "4. Klik 'Simpan'",
    "",
    "Driver dan Admin dapat mencatat penggunaan BBM."
])

# SLIDE 11: Laporan
add_content_slide(prs, "Langkah 8: Generate Laporan", [
    "1. Akses menu 'Laporan' atau 'Reports'",
    "2. Pilih jenis laporan:",
    "   • Laporan Riwayat Perjalanan",
    "   • Laporan Penggunaan BBM",
    "3. Filter berdasarkan:",
    "   • Tanggal",
    "   • Kendaraan",
    "   • Driver",
    "4. Klik 'Generate' atau 'Export PDF'",
    "",
    "Laporan dapat di-export ke format PDF."
])

# SLIDE 12: Penutup
add_title_slide(prs, "TERIMA KASIH", "Semoga prosedur ini membantu")

prs.save('E:/Laravel/pinjammobil/1-procedural.pptx')
print("Created 1-procedural.pptx")