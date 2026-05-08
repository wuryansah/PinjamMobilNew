from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 212, 255)
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 215, 0)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(12)
    return slide

# SLIDE 1: Title
add_title_slide(prs, "CARA PENGGUNAAN", "SISTEM PINJAM MOBIL")

# SLIDE 2: Login Admin
add_content_slide(prs, "Login sebagai ADMIN", [
    "1. Buka halaman web: domain.com/login",
    "2. Masukkan email admin",
    "3. Masukkan password",
    "4. Klik tombol 'Masuk' atau 'Login'",
    "",
    "Dashboard Admin:",
    "• Menu tambah kendaraan",
    "• Kelola driver",
    "• Setetujui permintaan",
    "• Laporan"
])

# SLIDE 3: Login Manager
add_content_slide(prs, "Login sebagai MANAGER", [
    "1. Buka halaman login",
    "2. Masukkan email manager",
    "3. Masukkan password",
    "4. Klik 'Masuk'",
    "",
    "Dashboard Manager:",
    "• Melihat dan menyetujui permintaan bawahannya"
])

# SLIDE 4: Login Driver
add_content_slide(prs, "Login sebagai DRIVER", [
    "1. Buka halaman login",
    "2. Masukkan email driver",
    "3. Masukkan password",
    "4. Klik 'Masuk'",
    "",
    "Dashboard Driver:",
    "• Melihat trip yang ditugaskan",
    "• Mulai/selesaikan perjalanan"
])

# SLIDE 5: Login Karyawan
add_content_slide(prs, "Login sebagai KARYAWAN", [
    "1. Buka halaman login",
    "2. Masukkan email karyawan",
    "3. Masukkan password",
    "4. Klik 'Masuk'",
    "",
    "Dashboard Karyawan:",
    "• Membuat permintaan kendaraan",
    "• Melihat status permintaan"
])

# SLIDE 6: Tambah Kendaraan
add_content_slide(prs, "Tambah Kendaraan (Admin)", [
    "1. Klik menu 'Kendaraan' di sidebar",
    "2. Klik tombol '+ Tambah'",
    "3. Isi data kendaraan:",
    "   • Plat nomor",
    "   • Merek dan tipe",
    "   • Jenis kendaraan",
    "   • Tahun pembuatan",
    "   • Kilometer saat ini",
    "   • Status ketersediaan",
    "4. Klik 'Simpan'",
    "",
    "Kendaraan baru siap ditugaskan."
])

# SLIDE 7: Buat Permintaan Karyawan
add_content_slide(prs, "Buat Permintaan (Karyawan)", [
    "1. Klik menu 'Permintaan'",
    "2. Klik tombol '+ Buat Permintaan'",
    "3. Isi formulir:",
    "   • Tujuan: Lokasi tujuan",
    "   • Keperluan: Untuk apa kendaraan",
    "   • Tanggal Mulai: Kapan membutuhkan",
    "   • Tanggal Selesai: Kapan mengembalikan",
    "   • Waktu Mulai & Selesai",
    "4. Klik 'Kirim Permintaan'"
])

# SLIDE 8: Approve Manager
add_content_slide(prs, "Setuju Permintaan (Manager)", [
    "1. Login sebagai Manager",
    "2. Lihat notifikasi permintaan baru",
    "3. Klik menu 'Permintaan'",
    "4. Pilih permintaan yang masuk",
    "5. Klik tombol 'Setuju' atau 'Tolak'",
    "   • Jika setuju: Klik tanpa catatan",
    "   • Jika tolak: Beri alasan",
    "",
    "Setelah di-setuju manager, request masuk ke admin."
])

# SLIDE 9: Assign Admin
add_content_slide(prs, "Assign Driver & Kendaraan (Admin)", [
    "1. Login sebagai Admin",
    "2. Pergi ke menu 'Permintaan'",
    "3. Cari request status 'Disetujui Manager'",
    "4. Klik 'Detail'",
    "5. Pilih kendaraan (yang tersedia)",
    "6. Pilih driver",
    "7. Tambah catatan jika perlu",
    "8. Klik 'Setuju & Tugaskan'"
])

# SLIDE 10: Mulai Perjalanan Driver
add_content_slide(prs, "Mulai Perjalanan (Driver)", [
    "1. Login sebagai Driver",
    "2. Lihat notifikasi penugasan baru",
    "3. Klik menu 'Permintaan'",
    "4. Cari trip yang ditugaskan",
    "5. Klik 'Mulai Perjalanan'",
    "6. Input Kilometer Awal (lihat odometer)",
    "7. Klik 'Mulai'",
    "",
    "Status berubah: Dalam Perjalanan"
])

# SLIDE 11: Selesai Perjalanan
add_content_slide(prs, "Selesai Perjalanan (Driver)", [
    "1. Setelah sampai dan kembali ke base",
    "2. Pergi ke menu 'Permintaan'",
    "3. Cari trip yang sedang berlangsung",
    "4. Klik 'Selesai Perjalanan'",
    "5. Input Kilometer Akhir (lihat odometer)",
    "6. Tambah catatan jika perlu",
    "7. Klik 'Selesai'",
    "",
    "Status menjadi Selesai. Kendaraan tersedia."
])

# SLIDE 12: Catat BBM
add_content_slide(prs, "Catat Penggunaan BBM", [
    "1. Klik menu 'BBM'",
    "2. Klik tombol '+ Tambah'",
    "3. Isi formulir:",
    "   • Pilih Kendaraan",
    "   • Tanggal pengisian",
    "   • Jumlah Liter",
    "   • Total Biaya",
    "   • Km Saat Ini",
    "   • Nota: Upload foto (opsional)",
    "4. Klik 'Simpan'"
])

# SLIDE 13: Lihat Laporan
add_content_slide(prs, "Lihat Laporan (Admin)", [
    "1. Login sebagai Admin",
    "2. Klik menu 'Laporan'",
    "3. Pilih jenis laporan:",
    "   • Riwayat Perjalanan",
    "   • Laporan BBM",
    "4. Atur filter:",
    "   • Tanggal mulai - selesai",
    "   • Kendaraan tertentu",
    "   • Driver tertentu",
    "5. Klik 'Tampilkan'",
    "6. Klik 'Export PDF' untuk download"
])

# SLIDE 14: Kelola User
add_content_slide(prs, "Kelola User (Admin)", [
    "1. Klik menu 'Pengguna'",
    "2. Lihat daftar semua user",
    "3. Tambah User Baru:",
    "   • Klik '+ Tambah'",
    "   • Isi: Nama, Email, Role, Dept",
    "   • Password default: 'password'",
    "4. Edit User:",
    "   • Klik tombol edit",
    "   • Ubah data",
    "   • Simpan",
    "5. Hapus User:",
    "   • Klik hapus",
    "   • Konfirmasi"
])

# SLIDE 15: Penutup
add_title_slide(prs, "TERIMA KASIH", "Sekarang Anda siap menggunakan sistem!")

prs.save('E:/Laravel/pinjammobil/2-usage.pptx')
print("Created 2-usage.pptx")